#!/usr/bin/env python3
"""
Generate PDF Documentation and PowerPoint Presentation
for the Inventory & POS System project.
"""

import os
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm, mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor, black, white
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
    PageBreak, ListFlowable, ListItem, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================
# PDF DOCUMENTATION
# ============================================================

def create_pdf():
    output_path = os.path.join(os.path.dirname(__file__), 'Dokumentasi_Arsitektur_Inventory_POS.pdf')
    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        rightMargin=2*cm,
        leftMargin=2*cm,
        topMargin=2*cm,
        bottomMargin=2*cm
    )

    styles = getSampleStyleSheet()
    
    # Custom styles
    styles.add(ParagraphStyle(
        'CoverTitle', parent=styles['Title'],
        fontSize=28, leading=34, spaceAfter=10,
        textColor=HexColor('#1a1a2e')
    ))
    styles.add(ParagraphStyle(
        'CoverSubtitle', parent=styles['Normal'],
        fontSize=14, leading=18, spaceAfter=30,
        textColor=HexColor('#4a4a6a'), alignment=TA_CENTER
    ))
    styles.add(ParagraphStyle(
        'SectionTitle', parent=styles['Heading1'],
        fontSize=18, leading=22, spaceBefore=20, spaceAfter=10,
        textColor=HexColor('#1a1a2e')
    ))
    styles.add(ParagraphStyle(
        'SubSection', parent=styles['Heading2'],
        fontSize=14, leading=17, spaceBefore=15, spaceAfter=8,
        textColor=HexColor('#2d2d5e')
    ))
    styles.add(ParagraphStyle(
        'BodyJustify', parent=styles['Normal'],
        fontSize=10, leading=14, spaceAfter=6,
        alignment=TA_JUSTIFY
    ))
    styles.add(ParagraphStyle(
        'BulletText', parent=styles['Normal'],
        fontSize=10, leading=14, leftIndent=20, spaceAfter=4
    ))
    styles.add(ParagraphStyle(
        'CodeStyle', parent=styles['Normal'],
        fontSize=9, leading=12, leftIndent=10,
        fontName='Courier', backColor=HexColor('#f5f5f5')
    ))
    styles.add(ParagraphStyle(
        'TableHeader', parent=styles['Normal'],
        fontSize=9, leading=12, textColor=white,
        alignment=TA_CENTER
    ))

    elements = []

    # --- COVER PAGE ---
    elements.append(Spacer(1, 6*cm))
    elements.append(Paragraph("DOKUMENTASI TEKNIS<br/>& ARSITEKTUR SISTEM", styles['CoverTitle']))
    elements.append(Spacer(1, 1*cm))
    elements.append(Paragraph("Inventory Management & Point of Sale System", styles['CoverSubtitle']))
    elements.append(Spacer(1, 2*cm))
    
    cover_data = [
        ['Versi', '1.0'],
        ['Tanggal', 'Agustus 2026'],
        ['Platform', 'Web Application'],
        ['Stack', 'Next.js + NestJS + PostgreSQL'],
    ]
    cover_table = Table(cover_data, colWidths=[4*cm, 8*cm])
    cover_table.setStyle(TableStyle([
        ('FONTSIZE', (0, 0), (-1, -1), 10),
        ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
        ('TOPPADDING', (0, 0), (-1, -1), 8),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#cccccc')),
    ]))
    elements.append(cover_table)
    elements.append(PageBreak())

    # --- TABLE OF CONTENTS ---
    elements.append(Paragraph("DAFTAR ISI", styles['SectionTitle']))
    toc_items = [
        "1. Executive Summary",
        "2. Arsitektur Sistem",
        "3. Technology Stack",
        "4. Database Schema",
        "5. Modul & Fitur",
        "6. Alur Bisnis (Business Flow)",
        "7. Security & Best Practices",
        "8. API Endpoints",
        "9. Deployment Guide",
        "10. Panduan Pengembangan",
    ]
    for item in toc_items:
        elements.append(Paragraph(item, styles['BodyJustify']))
    elements.append(PageBreak())

    # --- 1. EXECUTIVE SUMMARY ---
    elements.append(Paragraph("1. Executive Summary", styles['SectionTitle']))
    elements.append(Paragraph(
        "Sistem Inventory & POS ini dirancang sebagai solusi terintegrasi untuk manajemen gudang, "
        "distribusi barang ke outlet, dan proses penjualan (Point of Sale). Sistem mendukung "
        "multi-warehouse dan multi-outlet, dengan fitur tracking batch per produk menggunakan QR Code, "
        "manajemen expired date secara otomatis, serta integrasi payment gateway Midtrans.",
        styles['BodyJustify']
    ))
    elements.append(Spacer(1, 0.5*cm))
    elements.append(Paragraph("Keunggulan Utama:", styles['SubSection']))
    benefits = [
        "Dynamic RBAC — Role dan permission diatur admin, menu menyesuaikan otomatis",
        "QR Code Tracking — Setiap batch produk punya QR unik dengan data expired",
        "FIFO Expired Management — Batch terlama terjual duluan, expired otomatis terblokir",
        "Multi-location — Mendukung banyak gudang dan outlet",
        "Split Payment — Tunai, Midtrans, dan kombinasi keduanya",
        "Realtime Stock — Pergerakan stok tercatat di setiap transaksi",
        "Flexible Schema — Database dirancang generic untuk berbagai jenis toko",
    ]
    for b in benefits:
        elements.append(Paragraph(f"• {b}", styles['BulletText']))
    elements.append(PageBreak())

    # --- 2. ARSITEKTUR SISTEM ---
    elements.append(Paragraph("2. Arsitektur Sistem", styles['SectionTitle']))
    elements.append(Paragraph("Sistem menggunakan arsitektur client-server dengan pemisahan frontend dan backend:", styles['BodyJustify']))
    elements.append(Spacer(1, 0.5*cm))
    
    arch_data = [
        ['Layer', 'Technology', 'Deskripsi'],
        ['Frontend', 'Next.js 14 + React 18', 'Dashboard admin, POS interface, SSR/CSR'],
        ['Backend API', 'NestJS 10 + TypeScript', 'REST API, business logic, auth'],
        ['Database', 'PostgreSQL 16', 'Primary data store, relational'],
        ['Cache (opt)', 'Redis 7', 'Session, caching (future)'],
        ['Payment', 'Midtrans Snap', 'Payment gateway integration'],
    ]
    arch_table = Table(arch_data, colWidths=[3*cm, 4.5*cm, 9*cm])
    arch_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#1a1a2e')),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#cccccc')),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    elements.append(arch_table)
    elements.append(Spacer(1, 1*cm))

    elements.append(Paragraph("Arsitektur Diagram:", styles['SubSection']))
    elements.append(Paragraph(
        "[Browser/Client] → [Next.js Frontend :3001] → [NestJS API :3000] → [PostgreSQL :5432]",
        styles['CodeStyle']
    ))
    elements.append(Paragraph(
        "                                                    ↓",
        styles['CodeStyle']
    ))
    elements.append(Paragraph(
        "                                              [Midtrans API]",
        styles['CodeStyle']
    ))
    elements.append(PageBreak())

    # --- 3. TECHNOLOGY STACK ---
    elements.append(Paragraph("3. Technology Stack", styles['SectionTitle']))
    
    elements.append(Paragraph("Backend:", styles['SubSection']))
    be_items = [
        "NestJS 10 — Framework Node.js dengan arsitektur modular",
        "TypeORM — ORM untuk PostgreSQL dengan migration support",
        "Passport + JWT — Authentication dengan access/refresh token",
        "class-validator — Input validation dan sanitization",
        "Helmet — HTTP security headers",
        "@nestjs/throttler — Rate limiting",
        "@nestjs/schedule — Cron jobs (expired check)",
        "@nestjs/swagger — Auto-generated API documentation",
        "nestjs-pino — Structured logging",
    ]
    for item in be_items:
        elements.append(Paragraph(f"• {item}", styles['BulletText']))

    elements.append(Spacer(1, 0.5*cm))
    elements.append(Paragraph("Frontend:", styles['SubSection']))
    fe_items = [
        "Next.js 14 — React framework dengan App Router",
        "TypeScript — Type safety end-to-end",
        "Tailwind CSS 3 — Utility-first CSS framework",
        "Zustand — Lightweight state management",
        "Axios — HTTP client dengan interceptor",
        "Lucide React — Icon library",
        "Radix UI — Accessible headless components",
    ]
    for item in fe_items:
        elements.append(Paragraph(f"• {item}", styles['BulletText']))
    elements.append(PageBreak())

    # --- 4. DATABASE SCHEMA ---
    elements.append(Paragraph("4. Database Schema", styles['SectionTitle']))
    elements.append(Paragraph(
        "Database dirancang dengan prinsip normalisasi dan flexibility. "
        "Menggunakan UUID sebagai primary key, soft delete, dan timestamp tracking.",
        styles['BodyJustify']
    ))
    elements.append(Spacer(1, 0.5*cm))
    
    tables_info = [
        ['Entity/Table', 'Modul', 'Deskripsi'],
        ['users', 'Auth', 'Data user (email, password hash, profile)'],
        ['refresh_tokens', 'Auth', 'JWT refresh token storage'],
        ['roles', 'RBAC', 'Role definitions'],
        ['permissions', 'RBAC', 'Permission definitions (module:action)'],
        ['role_permissions', 'RBAC', 'Many-to-many role-permission'],
        ['user_roles', 'RBAC', 'User-role assignments'],
        ['modules', 'RBAC', 'Dynamic menu/module tree'],
        ['warehouses', 'Master', 'Data gudang'],
        ['outlets', 'Master', 'Data outlet/toko'],
        ['categories', 'Master', 'Kategori produk (tree structure)'],
        ['user_locations', 'Master', 'Assignment user ke lokasi'],
        ['products', 'Inventory', 'Master produk (SKU, harga, atribut)'],
        ['product_batches', 'Inventory', 'Batch per produk (expired, QR)'],
        ['inventory', 'Inventory', 'Stok per lokasi per batch'],
        ['stock_movements', 'Inventory', 'History pergerakan stok'],
        ['purchase_orders', 'PO', 'Surat pengiriman barang'],
        ['po_items', 'PO', 'Detail item dalam PO'],
        ['transactions', 'POS', 'Transaksi penjualan'],
        ['transaction_items', 'POS', 'Item dalam transaksi'],
        ['payments', 'POS', 'Pembayaran (cash/midtrans)'],
        ['members', 'POS', 'Data member/loyalty'],
        ['point_transactions', 'POS', 'History point member'],
        ['promos', 'POS', 'Promo & diskon engine'],
        ['returns', 'POS', 'Retur barang'],
        ['return_items', 'POS', 'Detail item retur'],
        ['notifications', 'System', 'Notifikasi user'],
    ]
    db_table = Table(tables_info, colWidths=[4*cm, 2.5*cm, 10*cm])
    db_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#1a1a2e')),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 8),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#cccccc')),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
        ('TOPPADDING', (0, 0), (-1, -1), 4),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f8fc')]),
    ]))
    elements.append(db_table)
    elements.append(PageBreak())

    # --- 5. MODUL & FITUR ---
    elements.append(Paragraph("5. Modul & Fitur", styles['SectionTitle']))
    
    modules = [
        ("Authentication", [
            "Register, Login, Logout",
            "JWT Access Token (15 menit) + Refresh Token (7 hari)",
            "Password hashing dengan bcrypt (cost factor 12)",
            "Rate limiting pada endpoint auth",
        ]),
        ("RBAC (Role-Based Access Control)", [
            "Dynamic roles — dibuat dan diatur admin",
            "Permission format: module:action (e.g. inventory:read)",
            "Guard decorator @RequirePermission('...')",
            "Menu dashboard otomatis menyesuaikan permission user",
        ]),
        ("Master Data", [
            "CRUD Warehouses, Outlets, Categories",
            "Category tree structure (parent-child)",
            "User-location assignment (user hanya akses lokasi yang di-assign)",
        ]),
        ("Inventory & Product", [
            "CRUD Products dengan dynamic attributes (JSONB)",
            "Batch management (production date, expired date)",
            "QR Code generation per batch",
            "Stock tracking per lokasi per batch",
            "Stock movements logging",
            "Cron job daily: auto-block batch expired",
        ]),
        ("Purchase Order (PO)", [
            "Flow: Draft → Submitted → Shipped → Received → Approved",
            "Stock otomatis terpotong dari warehouse saat ship",
            "Stock otomatis bertambah di outlet saat approve",
            "Search by PO number (untuk scan/input di outlet)",
        ]),
        ("POS (Point of Sale)", [
            "Scan QR / search product → add to cart",
            "FIFO: batch paling dekat expired dijual duluan",
            "Block otomatis: batch expired tidak bisa dijual",
            "Diskon per item & per transaksi",
            "Hold & recall transaction",
            "Void transaction",
        ]),
        ("Payment", [
            "Cash payment (hitung kembalian)",
            "Midtrans integration (Snap token)",
            "Split payment (sebagian cash, sebagian digital)",
            "Webhook untuk notification dari Midtrans",
        ]),
        ("Member & Loyalty", [
            "Registrasi member (name, phone, email)",
            "Point earning: 1 point per Rp 10.000",
            "Tier system: Bronze, Silver, Gold, Platinum",
            "Point redemption (future enhancement)",
        ]),
        ("Reports", [
            "Sales report (harian/mingguan/bulanan/per outlet)",
            "Stock report & low stock alert",
            "Expired/waste report",
            "Profit/Loss per outlet",
            "Cashier recap",
            "Member report",
            "Export to Excel/PDF",
        ]),
        ("Notifications", [
            "Expired warning (H-7, H-3, H-1)",
            "Low stock alert",
            "PO status change notification",
        ]),
    ]

    for mod_name, features in modules:
        elements.append(Paragraph(mod_name, styles['SubSection']))
        for f in features:
            elements.append(Paragraph(f"• {f}", styles['BulletText']))
        elements.append(Spacer(1, 0.3*cm))
    elements.append(PageBreak())

    # --- 6. BUSINESS FLOW ---
    elements.append(Paragraph("6. Alur Bisnis (Business Flow)", styles['SectionTitle']))
    
    elements.append(Paragraph("A. Alur Input Barang & Distribusi:", styles['SubSection']))
    flow_a = [
        "1. Staff gudang input produk baru ke sistem",
        "2. Buat batch (set tanggal produksi & expired)",
        "3. Sistem generate QR Code → print & tempel di barang",
        "4. Stock In ke warehouse (stok bertambah)",
        "5. Buat Purchase Order (pilih outlet tujuan, pilih batch)",
        "6. Submit PO → Review → Ship (stok gudang berkurang)",
        "7. Outlet scan/input nomor PO → detail muncul",
        "8. Outlet verifikasi & approve → stok outlet bertambah",
    ]
    for f in flow_a:
        elements.append(Paragraph(f, styles['BulletText']))

    elements.append(Spacer(1, 0.5*cm))
    elements.append(Paragraph("B. Alur Penjualan (POS):", styles['SubSection']))
    flow_b = [
        "1. Kasir scan QR produk atau search manual",
        "2. Sistem lookup batch (FIFO) + check expired",
        "3. Produk masuk ke cart → hitung subtotal",
        "4. Input member (optional) → apply promo/diskon",
        "5. Pilih metode bayar: Cash / Midtrans / Split",
        "6. Proses pembayaran → stok berkurang → movement tercatat",
        "7. Member dapat point (jika ada)",
        "8. Cetak struk / kirim digital receipt",
    ]
    for f in flow_b:
        elements.append(Paragraph(f, styles['BulletText']))
    elements.append(PageBreak())

    # --- 7. SECURITY ---
    elements.append(Paragraph("7. Security & Best Practices", styles['SectionTitle']))
    security_items = [
        ("Helmet", "HTTP security headers (CSP, HSTS, X-Frame-Options, dll)"),
        ("Rate Limiting", "100 req/menit per IP (configurable)"),
        ("CORS", "Whitelist origin frontend only"),
        ("JWT", "Short-lived access token + long-lived refresh token"),
        ("Bcrypt", "Password hashing dengan cost factor 12"),
        ("Validation", "class-validator whitelist mode — reject unknown fields"),
        ("Soft Delete", "Data tidak pernah benar-benar dihapus"),
        ("SQL Injection", "TypeORM parameterized queries"),
        ("XSS", "React auto-escaping + CSP headers"),
        ("Audit Trail", "Stock movements, created_by fields"),
    ]
    for title, desc in security_items:
        elements.append(Paragraph(f"<b>{title}</b> — {desc}", styles['BulletText']))
    elements.append(PageBreak())

    # --- 8. API ENDPOINTS ---
    elements.append(Paragraph("8. API Endpoints", styles['SectionTitle']))
    elements.append(Paragraph("Base URL: http://localhost:3000/api/v1", styles['CodeStyle']))
    elements.append(Paragraph("Swagger Docs: http://localhost:3000/docs", styles['CodeStyle']))
    elements.append(Spacer(1, 0.5*cm))

    endpoints = [
        ['Method', 'Endpoint', 'Description'],
        ['POST', '/auth/register', 'Register user baru'],
        ['POST', '/auth/login', 'Login → access + refresh token'],
        ['POST', '/auth/refresh', 'Refresh access token'],
        ['POST', '/auth/logout', 'Revoke refresh token'],
        ['GET', '/rbac/roles', 'List semua roles'],
        ['POST', '/rbac/roles', 'Create role'],
        ['GET', '/rbac/permissions', 'List permissions'],
        ['GET', '/rbac/menu', 'Get menu sesuai user permission'],
        ['CRUD', '/master/warehouses', 'Manage warehouses'],
        ['CRUD', '/master/outlets', 'Manage outlets'],
        ['CRUD', '/master/categories', 'Manage categories'],
        ['CRUD', '/inventory/products', 'Manage products'],
        ['POST', '/inventory/batches', 'Create batch + generate QR'],
        ['POST', '/inventory/stock-in', 'Stock in ke warehouse'],
        ['GET', '/inventory/stock/:type/:id', 'Get stock by location'],
        ['POST', '/purchase-orders', 'Create PO'],
        ['PUT', '/purchase-orders/:id/ship', 'Ship PO'],
        ['PUT', '/purchase-orders/:id/receive', 'Receive PO'],
        ['PUT', '/purchase-orders/:id/approve', 'Approve PO'],
        ['POST', '/pos/transactions', 'Create transaction'],
        ['POST', '/pos/payment', 'Process payment'],
        ['POST', '/pos/members', 'Register member'],
        ['GET', '/notifications', 'Get user notifications'],
    ]
    ep_table = Table(endpoints, colWidths=[2*cm, 5.5*cm, 9*cm])
    ep_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#1a1a2e')),
        ('TEXTCOLOR', (0, 0), (-1, 0), white),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 8),
        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#cccccc')),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
        ('TOPPADDING', (0, 0), (-1, -1), 3),
        ('FONTNAME', (1, 1), (1, -1), 'Courier'),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f8fc')]),
    ]))
    elements.append(ep_table)
    elements.append(PageBreak())

    # --- 9. DEPLOYMENT ---
    elements.append(Paragraph("9. Deployment Guide", styles['SectionTitle']))
    elements.append(Paragraph("Development:", styles['SubSection']))
    dev_steps = [
        "1. docker compose up -d  (PostgreSQL + Redis)",
        "2. cd be && cp .env.example .env && npm install",
        "3. npm run start:dev  (Backend di port 3000)",
        "4. npm run seed  (Seed admin user & permissions)",
        "5. cd fe && npm install && npm run dev  (Frontend di port 3001)",
        "6. Buka http://localhost:3001 → Login: admin@inventorypos.com / Admin@123",
    ]
    for s in dev_steps:
        elements.append(Paragraph(s, styles['BulletText']))
    
    elements.append(Spacer(1, 0.5*cm))
    elements.append(Paragraph("Production:", styles['SubSection']))
    prod_steps = [
        "• Set NODE_ENV=production di .env",
        "• Ganti JWT_SECRET dan JWT_REFRESH_SECRET dengan random secure key",
        "• Set DB credentials production",
        "• Frontend: npm run build → deploy ke Vercel/server",
        "• Backend: npm run build → pm2 start dist/main.js",
        "• Setup reverse proxy (Nginx) dengan HTTPS",
        "• Disable synchronize di TypeORM, gunakan migrations",
    ]
    for s in prod_steps:
        elements.append(Paragraph(s, styles['BulletText']))
    elements.append(PageBreak())

    # --- 10. PANDUAN PENGEMBANGAN ---
    elements.append(Paragraph("10. Panduan Pengembangan", styles['SectionTitle']))
    elements.append(Paragraph("Struktur Project:", styles['SubSection']))
    structure = """
    inventory-pos/
    ├── docker-compose.yml
    ├── be/                    (NestJS Backend)
    │   ├── src/
    │   │   ├── main.ts
    │   │   ├── app.module.ts
    │   │   ├── config/        (DB, env validation)
    │   │   ├── common/        (BaseEntity, DTOs, filters)
    │   │   ├── health/        (Health check)
    │   │   └── modules/
    │   │       ├── auth/      (JWT, login, register)
    │   │       ├── rbac/      (Roles, permissions, menu)
    │   │       ├── master/    (Warehouse, outlet, category)
    │   │       ├── inventory/ (Product, batch, stock)
    │   │       ├── purchase-order/
    │   │       ├── pos/       (Transaction, payment, member)
    │   │       └── notifications/
    │   └── package.json
    └── fe/                    (Next.js Frontend)
        ├── src/
        │   ├── app/           (Pages - App Router)
        │   ├── components/    (UI components)
        │   ├── lib/           (API client, utils)
        │   ├── store/         (Zustand stores)
        │   └── types/         (TypeScript interfaces)
        └── package.json
    """
    for line in structure.strip().split('\n'):
        elements.append(Paragraph(line, styles['CodeStyle']))

    # Build PDF
    doc.build(elements)
    print(f"PDF created: {output_path}")
    return output_path


# ============================================================
# POWERPOINT PRESENTATION
# ============================================================

def add_slide(prs, title, content_items, layout_idx=1):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(8)
    
    return slide


def add_title_slide(prs, title, subtitle):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    
    slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[2]  # Section header
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    return slide


def create_pptx():
    output_path = os.path.join(os.path.dirname(__file__), 'Proposal_Pengadaan_Inventory_POS.pptx')
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    add_title_slide(prs, 
        "Proposal Pengadaan\nSistem Inventory & POS",
        "Solusi Terintegrasi Manajemen Gudang, Distribusi & Penjualan\nAgustus 2026"
    )

    # Slide 2: Latar Belakang
    add_slide(prs, "Latar Belakang", [
        "Pengelolaan stok manual rawan kesalahan & kehilangan data",
        "Tracking expired date sulit dilakukan secara konsisten",
        "Distribusi gudang ke outlet tanpa sistem terintegrasi",
        "Tidak ada visibilitas realtime terhadap stok & penjualan",
        "Kebutuhan payment digital (cashless) meningkat",
        "Pelaporan manual memakan waktu & tidak akurat",
    ])

    # Slide 3: Tujuan
    add_slide(prs, "Tujuan Project", [
        "Membangun sistem inventory terintegrasi gudang-outlet",
        "Automasi tracking expired date & blocking produk kadaluarsa",
        "Digitalisasi alur distribusi dengan Purchase Order system",
        "POS modern dengan QR scan, split payment, & loyalty",
        "Dashboard reporting realtime untuk pengambilan keputusan",
        "Sistem scalable — bisa dipakai untuk berbagai jenis usaha",
    ])

    # Slide 4: Fitur Utama
    add_slide(prs, "Fitur Utama", [
        "Dynamic RBAC — Role & akses diatur fleksibel oleh admin",
        "QR Code per Batch — Scan untuk info produk & expired date",
        "Purchase Order — Alur pengiriman gudang → outlet lengkap",
        "POS Terminal — Fast, user-friendly, support split payment",
        "Midtrans Integration — Pembayaran digital (QRIS, e-wallet, dll)",
        "Member & Loyalty — Point system, tier membership",
        "Auto Expired Block — Barang kadaluarsa otomatis tidak bisa dijual",
        "Reports — Sales, stock, profit/loss, expired/waste, cashier recap",
    ])

    # Slide 5: Arsitektur
    add_slide(prs, "Arsitektur Sistem", [
        "Frontend: Next.js 14 (React) — Dashboard & POS UI",
        "Backend: NestJS (Node.js) — REST API, Business Logic",
        "Database: PostgreSQL 16 — Reliable, scalable RDBMS",
        "Auth: JWT (Access + Refresh Token) — Secure & stateless",
        "Payment: Midtrans Snap — PCI-DSS compliant gateway",
        "Deployment: Docker ready, CI/CD compatible",
    ])

    # Slide 6: Alur Bisnis
    add_slide(prs, "Alur Bisnis", [
        "1. Input Barang → Generate QR Code → Tempel di produk",
        "2. Stock In ke Gudang → Stok tercatat per batch",
        "3. Buat PO → Kirim ke Outlet → Outlet Terima & Approve",
        "4. Stok otomatis berpindah dari gudang ke outlet",
        "5. Kasir scan QR / search → Bayar (Cash/Digital)",
        "6. Stok berkurang otomatis, laporan ter-update realtime",
    ])

    # Slide 7: Keamanan
    add_slide(prs, "Keamanan Sistem", [
        "JWT Authentication — Token-based, short-lived",
        "Rate Limiting — Proteksi terhadap brute force & DDoS",
        "Helmet — Security headers (CSP, HSTS, X-Frame, dll)",
        "Input Validation — Whitelist mode, reject unknown fields",
        "Bcrypt — Password hashing (industry standard)",
        "Soft Delete — Data tidak pernah benar-benar hilang",
        "CORS — Hanya frontend terdaftar yang bisa akses API",
        "Captcha — Proteksi form login dari bot",
    ])

    # Slide 8: Modul Sistem
    add_slide(prs, "Modul Sistem", [
        "Auth & RBAC — Login, register, dynamic roles & permissions",
        "Master Data — Warehouse, outlet, category management",
        "Inventory — Product, batch, stock tracking, QR code",
        "Purchase Order — Distribution flow warehouse → outlet",
        "POS — Transaction, cart, discount, FIFO expired",
        "Payment — Cash, Midtrans, split payment",
        "Member — Loyalty, points, tiers",
        "Reports — Sales, stock, expired, profit/loss, cashier",
        "Notifications — Expired alert, low stock, PO status",
    ])

    # Slide 9: Timeline
    add_slide(prs, "Timeline Pengembangan", [
        "Minggu 1-2: Setup, Auth, RBAC, Master Data",
        "Minggu 3-4: Inventory, Product, Batch, QR Code",
        "Minggu 5-6: Purchase Order, Stock Movement",
        "Minggu 7-8: POS, Payment, Member/Loyalty",
        "Minggu 9-10: Reports, Notifications, Promo Engine",
        "Minggu 11-12: Frontend UI (Dashboard, POS, dll)",
        "Minggu 13-14: Integration Testing & Security Audit",
        "Minggu 15-16: UAT, Bug Fix, Deployment",
    ])

    # Slide 10: Tech Stack
    add_slide(prs, "Technology Stack", [
        "Frontend: Next.js 14, React 18, TypeScript, Tailwind CSS",
        "Backend: NestJS 10, TypeORM, PostgreSQL 16",
        "Auth: JWT + Passport, bcrypt, Refresh Token",
        "Payment: Midtrans Snap API",
        "DevOps: Docker, Docker Compose",
        "Tools: Swagger (API docs), Pino (logging)",
    ])

    # Slide 11: Deliverables
    add_slide(prs, "Deliverables", [
        "Source code (Backend + Frontend) — fully documented",
        "Database schema & migration scripts",
        "API Documentation (Swagger interactive docs)",
        "User manual & admin guide",
        "Deployment guide (dev & production)",
        "Training session untuk admin & staff",
    ])

    # Slide 12: Penutup
    add_title_slide(prs,
        "Terima Kasih",
        "Inventory & POS System — Solusi Terintegrasi\nuntuk Manajemen Gudang, Distribusi & Penjualan"
    )

    prs.save(output_path)
    print(f"PPTX created: {output_path}")
    return output_path


# ============================================================
# MAIN
# ============================================================

if __name__ == '__main__':
    print("Generating documentation...")
    pdf_path = create_pdf()
    pptx_path = create_pptx()
    print(f"\nDone!")
    print(f"  PDF:  {pdf_path}")
    print(f"  PPTX: {pptx_path}")
